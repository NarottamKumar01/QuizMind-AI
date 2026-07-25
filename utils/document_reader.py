from pypdf import PdfReader
from docx import Document
from pptx import Presentation


# ----------------------------
# Extract Text From PDF
# ----------------------------

def extract_pdf(path):

    try:

        reader = PdfReader(path)

        text = ""

        for page in reader.pages:

            page_text = page.extract_text()

            if page_text:

                text += page_text + "\n"


        return text.strip()


    except Exception as e:

        return f"Error reading PDF file: {str(e)}"



# ----------------------------
# Extract Text From DOCX
# ----------------------------

def extract_docx(path):

    try:

        document = Document(path)

        text = ""

        for paragraph in document.paragraphs:

            if paragraph.text.strip():

                text += paragraph.text + "\n"


        return text.strip()


    except Exception as e:

        return f"Error reading DOCX file: {str(e)}"



# ----------------------------
# Extract Text From PPTX
# ----------------------------

def extract_pptx(path):

    try:

        presentation = Presentation(path)

        text = ""

        for slide in presentation.slides:

            for shape in slide.shapes:

                if hasattr(shape, "text"):

                    if shape.text.strip():

                        text += shape.text + "\n"


        return text.strip()


    except Exception as e:

        return f"Error reading PPTX file: {str(e)}"



# ----------------------------
# Universal Document Reader
# ----------------------------

def extract_document(path):

    extension = path.lower().split(".")[-1]


    if extension == "pdf":

        return extract_pdf(path)


    elif extension == "docx":

        return extract_docx(path)


    elif extension == "pptx":

        return extract_pptx(path)


    else:

        return "Unsupported document format"